from __future__ import annotations

from getpass import getuser
from io import BytesIO
from pathlib import Path

from .fileops import FileService
from .workspace import Workspace


class AutomationDependencyError(RuntimeError):
    pass


def _deps():
    try:
        from docx import Document
        from openpyxl import Workbook, load_workbook
        from openpyxl.utils.cell import range_boundaries
        from pptx import Presentation
    except ImportError as exc:
        raise AutomationDependencyError(
            "Install pinned deterministic automation dependencies: pip install -r requirements-automation.txt"
        ) from exc
    return Document, Workbook, load_workbook, range_boundaries, Presentation


def sheet_range_to_document_table(workspace: Workspace, source_rel: str, sheet_name: str, cell_range: str, dest_rel: str) -> Path:
    Document, _Workbook, load_workbook, range_boundaries, _Presentation = _deps()
    source = workspace.paths.resolve(source_rel, must_exist=True)
    if source.suffix.lower() not in {".xlsx", ".xlsm"}:
        raise ValueError("Source must be XLSX/XLSM.")
    wb = load_workbook(source, data_only=False, read_only=True)
    try:
        if sheet_name not in wb.sheetnames:
            raise KeyError(f"Worksheet not found: {sheet_name}")
        ws = wb[sheet_name]
        min_col, min_row, max_col, max_row = range_boundaries(cell_range)
        rows = [[ws.cell(r, c).value for c in range(min_col, max_col + 1)] for r in range(min_row, max_row + 1)]
    finally:
        # read_only workbooks keep the package stream open on Windows until
        # explicitly closed, which otherwise prevents workspace cleanup.
        wb.close()
    if not rows:
        raise ValueError("Range is empty.")
    doc = Document()
    doc.add_heading("Sheet Range Export", level=1)
    doc.add_paragraph(f"Source: {source_rel} | Sheet: {sheet_name} | Range: {cell_range}")
    table = doc.add_table(rows=len(rows), cols=len(rows[0])); table.style = "Table Grid"
    for r, values in enumerate(rows):
        for c, value in enumerate(values):
            table.cell(r, c).text = "" if value is None else str(value)
    buf = BytesIO(); doc.save(buf)
    target = FileService(workspace).atomic_write(dest_rel, buf.getvalue(), preserve_version=True)
    workspace.audit.append(
        who=getuser(), what="workflow.sheet_range_to_document_table", target=dest_rel,
        intent=f"{source_rel}:{sheet_name}!{cell_range}", approval="USER", execution="DETERMINISTIC",
        result="CREATED", verification="DOCX written atomically",
        metadata={"source": source_rel, "sheet": sheet_name, "range": cell_range, "rows": len(rows), "columns": len(rows[0])},
    )
    return target


def document_table_to_sheet(workspace: Workspace, source_rel: str, table_index: int, dest_rel: str) -> Path:
    Document, Workbook, _load_workbook, _range_boundaries, _Presentation = _deps()
    source = workspace.paths.resolve(source_rel, must_exist=True)
    if source.suffix.lower() != ".docx":
        raise ValueError("Source must be DOCX.")
    doc = Document(source)
    if table_index < 0 or table_index >= len(doc.tables):
        raise IndexError("Document table index is out of range.")
    table = doc.tables[table_index]
    wb = Workbook(); ws = wb.active; ws.title = "ImportedTable"
    for row in table.rows:
        ws.append([cell.text for cell in row.cells])
    ws.freeze_panes = "A2" if len(table.rows) > 1 else None
    buf = BytesIO()
    try:
        wb.save(buf)
    finally:
        wb.close()
    target = FileService(workspace).atomic_write(dest_rel, buf.getvalue(), preserve_version=True)
    workspace.audit.append(
        who=getuser(), what="workflow.document_table_to_sheet", target=dest_rel,
        intent=f"{source_rel}:table[{table_index}]", approval="USER", execution="DETERMINISTIC",
        result="CREATED", verification="XLSX written atomically",
        metadata={"source": source_rel, "table_index": table_index, "rows": len(table.rows)},
    )
    return target


def document_outline_to_slides(workspace: Workspace, source_rel: str, dest_rel: str) -> Path:
    Document, _Workbook, _load_workbook, _range_boundaries, Presentation = _deps()
    source = workspace.paths.resolve(source_rel, must_exist=True)
    if source.suffix.lower() != ".docx":
        raise ValueError("Source must be DOCX.")
    doc = Document(source)
    headings: list[tuple[int, str]] = []
    for p in doc.paragraphs:
        name = p.style.name if p.style else ""
        if name.startswith("Heading") and p.text.strip():
            try:
                level = int(name.split()[-1])
            except ValueError:
                level = 1
            headings.append((level, p.text.strip()))
    if not headings:
        raise ValueError("No heading paragraphs found.")
    prs = Presentation()
    # Remove package default first slide only if created by template; Presentation() starts with zero slides.
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = headings[0][1]
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = f"Draft from {source_rel}"
    current_slide = None
    bullets = None
    for level, text in headings[1:]:
        if level <= 2 or current_slide is None:
            current_slide = prs.slides.add_slide(prs.slide_layouts[1])
            current_slide.shapes.title.text = text
            bullets = current_slide.placeholders[1].text_frame
            bullets.clear()
        else:
            p = bullets.add_paragraph() if bullets is not None else None
            if p is not None:
                p.text = text; p.level = min(level - 3, 4)
    buf = BytesIO(); prs.save(buf)
    target = FileService(workspace).atomic_write(dest_rel, buf.getvalue(), preserve_version=True)
    workspace.audit.append(
        who=getuser(), what="workflow.document_outline_to_slides", target=dest_rel,
        intent=source_rel, approval="USER", execution="DETERMINISTIC", result="CREATED",
        verification="PPTX written atomically", metadata={"source": source_rel, "heading_count": len(headings)},
    )
    return target
