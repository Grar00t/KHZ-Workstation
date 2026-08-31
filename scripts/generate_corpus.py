from __future__ import annotations

import os
import sys
from datetime import date, timedelta
from pathlib import Path

from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.comments import Comment
from openpyxl.formatting.rule import CellIsRule, FormulaRule
from openpyxl.styles import Font, PatternFill
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement as PPTXOxmlElement
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches as PInches
from pptx.util import Pt as PPt
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from PIL import Image, ImageDraw

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "acceptance" / "corpus"
OUT.mkdir(parents=True, exist_ok=True)


def tiny_png(path: Path) -> None:
    image = Image.new("RGB", (640, 320), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((30, 30, 610, 290), outline="black", width=3)
    draw.text((70, 125), "SYNTHETIC TEST DATA - KHZ", fill="black")
    image.save(path, format="PNG")


def add_page_number(paragraph):
    run = paragraph.add_run()
    fld = OxmlElement("w:fldSimple")
    fld.set(qn("w:instr"), "PAGE")
    run._r.addnext(fld)


def add_hyperlink(paragraph, text: str, url: str):
    part = paragraph.part
    rid = part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink", is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), rid)
    new_run = OxmlElement("w:r")
    rpr = OxmlElement("w:rPr")
    color = OxmlElement("w:color"); color.set(qn("w:val"), "0563C1")
    underline = OxmlElement("w:u"); underline.set(qn("w:val"), "single")
    rpr.append(color); rpr.append(underline)
    new_run.append(rpr)
    t = OxmlElement("w:t"); t.text = text
    new_run.append(t); hyperlink.append(new_run); paragraph._p.append(hyperlink)


def add_toc(paragraph):
    fld = OxmlElement("w:fldSimple")
    fld.set(qn("w:instr"), 'TOC \\o "1-3" \\h \\z \\u')
    paragraph._p.append(fld)


def add_tracked_insertion(paragraph, text: str):
    ins = OxmlElement("w:ins")
    ins.set(qn("w:author"), "KHZ Synthetic Fixture")
    ins.set(qn("w:date"), "2026-08-31T12:00:00Z")
    run = OxmlElement("w:r")
    t = OxmlElement("w:t"); t.text = text
    run.append(t); ins.append(run); paragraph._p.append(ins)


def make_docx(img: Path) -> Path:
    path = OUT / "InstitutionalReport.docx"
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.8); sec.bottom_margin = Inches(0.8); sec.left_margin = Inches(0.9); sec.right_margin = Inches(0.9)
    h = sec.header.paragraphs[0]; h.text = "SYNTHETIC TEST DATA | KHZ Institutional Report"; h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    f = sec.footer.paragraphs[0]; f.text = "Confidential fixture - Page "; add_page_number(f)
    title = doc.add_paragraph(); title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("Institutional Operations Report"); r.bold = True; r.font.size = Pt(24)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.add_run("SYNTHETIC TEST DATA").bold = True
    doc.add_page_break()
    doc.add_heading("Table of Contents", level=1); add_toc(doc.add_paragraph())
    doc.add_heading("1. Executive Summary", level=1)
    para = doc.add_paragraph("This synthetic report exercises common institutional document structures. ")
    add_hyperlink(para, "KHZ local documentation", "https://example.invalid/khz")
    doc.add_heading("1.1 Scope", level=2)
    p = doc.add_paragraph("The fixture contains styles, tables, images, headers, footers, lists, comments, tracked changes, and fields.")
    add_tracked_insertion(p, " Inserted synthetic tracked-change text.")
    doc.add_paragraph("Operational continuity", style="List Bullet")
    doc.add_paragraph("Evidence-based verification", style="List Bullet")
    doc.add_paragraph("Create backup", style="List Number")
    doc.add_paragraph("Validate backup", style="List Number")
    doc.add_heading("2. Metrics", level=1)
    table = doc.add_table(rows=1, cols=4); table.style = "Table Grid"
    headers = ["Department", "Cases", "Average Time", "Status"]
    for i, x in enumerate(headers): table.rows[0].cells[i].text = x
    for row in [("Operations", "124", "3.4", "On target"), ("Finance", "87", "2.1", "On target"), ("Quality", "33", "5.0", "Review")]:
        cells = table.add_row().cells
        for i, x in enumerate(row): cells[i].text = x
    doc.add_picture(str(img), width=Inches(1.0))
    cap = doc.add_paragraph("Figure 1. Synthetic fixture image"); cap.style = doc.styles["Caption"]
    doc.add_heading("3. Review", level=1)
    review = doc.add_paragraph("Review note target sentence for comment testing.")
    if hasattr(doc, "add_comment"):
        try:
            doc.add_comment(review.runs, text="Synthetic reviewer comment.", author="KHZ QA", initials="KQ")
        except Exception:
            pass
    doc.add_section().orientation = WD_ORIENT.PORTRAIT
    doc.save(path)
    return path


def make_xlsx() -> Path:
    path = OUT / "InstitutionalWorkbook.xlsx"
    wb = Workbook()
    summary = wb.active; summary.title = "Summary"
    tx = wb.create_sheet("Transactions")
    deps = wb.create_sheet("Departments")
    lookup = wb.create_sheet("Lookup")
    notes = wb.create_sheet("Notes")
    deps.append(["DepartmentCode", "DepartmentName"])
    for row in [("OPS", "Operations"), ("FIN", "Finance"), ("QLT", "Quality"), ("HR", "Human Resources"), ("IT", "Technology")]: deps.append(row)
    lookup.append(["Category", "TaxRate"])
    for row in [("Services", 0.05), ("Supplies", 0.10), ("Capital", 0.15)]: lookup.append(row)
    headers = ["ID", "Date", "Department", "Category", "Amount", "TaxRate", "Tax", "Total", "Approved", "DeptName"]
    tx.append(headers)
    start = date(2025, 1, 1)
    dep_codes = ["OPS", "FIN", "QLT", "HR", "IT"]
    cats = ["Services", "Supplies", "Capital"]
    for i in range(1, 1201):
        row = i + 1
        amount = 100 + (i * 37) % 5000
        tx.cell(row, 1, f"TX-{i:05d}")
        tx.cell(row, 2, start + timedelta(days=i % 365)); tx.cell(row, 2).number_format = "yyyy-mm-dd"
        tx.cell(row, 3, dep_codes[i % len(dep_codes)])
        tx.cell(row, 4, cats[i % len(cats)])
        tx.cell(row, 5, amount); tx.cell(row, 5).number_format = '$#,##0.00'
        tx.cell(row, 6, f'=VLOOKUP(D{row},Lookup!$A$2:$B$4,2,FALSE)'); tx.cell(row, 6).number_format = "0%"
        tx.cell(row, 7, f'=ROUND(E{row}*F{row},2)'); tx.cell(row, 7).number_format = '$#,##0.00'
        tx.cell(row, 8, f'=E{row}+G{row}'); tx.cell(row, 8).number_format = '$#,##0.00'
        tx.cell(row, 9, "Yes" if i % 7 else "No")
        tx.cell(row, 10, f'=_xlfn.XLOOKUP(C{row},Departments!$A$2:$A$6,Departments!$B$2:$B$6,"Unknown")')
    tab = Table(displayName="TransactionsTable", ref="A1:J1201")
    tab.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showRowStripes=True, showFirstColumn=False, showLastColumn=False)
    tx.add_table(tab)
    tx.freeze_panes = "A2"
    tx.auto_filter.ref = "A1:J1201"
    dv = DataValidation(type="list", formula1='"Services,Supplies,Capital"', allow_blank=False)
    tx.add_data_validation(dv); dv.add("D2:D1201")
    tx.conditional_formatting.add("H2:H1201", CellIsRule(operator="greaterThan", formula=["3000"], fill=PatternFill("solid", fgColor="FFF2CC")))
    tx.conditional_formatting.add("I2:I1201", FormulaRule(formula=['I2="No"'], font=Font(color="9C0006"), fill=PatternFill("solid", fgColor="FFC7CE")))
    tx["A2"].comment = Comment("SYNTHETIC TEST DATA comment", "KHZ QA")
    summary.append(["Metric", "Value"])
    summary.append(["Operations Total", '=SUMIFS(Transactions!$H$2:$H$1201,Transactions!$C$2:$C$1201,"OPS")'])
    summary.append(["Approved Average", '=AVERAGEIFS(Transactions!$H$2:$H$1201,Transactions!$I$2:$I$1201,"Yes")'])
    summary.append(["Rejected Count", '=COUNTIFS(Transactions!$I$2:$I$1201,"No")'])
    summary.append(["Maximum", '=MAX(Transactions!$H$2:$H$1201)'])
    summary.append(["Minimum", '=MIN(Transactions!$H$2:$H$1201)'])
    summary.append(["Lookup Example", '=_xlfn.XLOOKUP("OPS",Departments!$A$2:$A$6,Departments!$B$2:$B$6,"Unknown")'])
    summary.append(["Dynamic Array Example", '=_xlfn.UNIQUE(Transactions!$C$2:$C$1201)'])
    for cell in summary[1]: cell.font = Font(bold=True)
    summary.freeze_panes = "A2"
    summary.sheet_properties.pageSetUpPr.fitToPage = True
    summary.protection.sheet = True; summary.protection.set_password("synthetic")
    from openpyxl.workbook.defined_name import DefinedName
    wb.defined_names.add(DefinedName("TransactionAmounts", attr_text="Transactions!$H$2:$H$1201"))
    chart = BarChart(); chart.type = "col"; chart.title = "Synthetic Institutional Metrics"; chart.y_axis.title = "Value"; chart.x_axis.title = "Metric"
    data = Reference(summary, min_col=2, min_row=2, max_row=6); cats_ref = Reference(summary, min_col=1, min_row=2, max_row=6)
    chart.add_data(data, titles_from_data=False); chart.set_categories(cats_ref); chart.height = 7; chart.width = 12
    summary.add_chart(chart, "D2")
    notes["A1"] = "SYNTHETIC TEST DATA"; notes["A2"] = "Pivot table is added by the LibreOffice UNO spike when available."
    for ws in wb.worksheets:
        ws.sheet_view.showGridLines = True
    wb.save(path)
    return path


def add_transition(slide):
    from lxml.etree import QName
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    root = slide._element
    trans = PPTXOxmlElement("p:transition")
    trans.set("spd", "med")
    trans.append(PPTXOxmlElement("p:fade"))
    idx = len(root)
    for i, child in enumerate(root):
        if QName(child).localname in {"timing", "extLst"}:
            idx = i
            break
    root.insert(idx, trans)


def make_pptx(img: Path) -> Path:
    path = OUT / "InstitutionalPresentation.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0]); slide.shapes.title.text = "Institutional Review"; slide.placeholders[1].text = "SYNTHETIC TEST DATA"
    slide = prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text = "Agenda"; slide.placeholders[1].text = "Operations\nFinance\nQuality\nNext steps"
    slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = "Local-first evidence"; slide.shapes.add_picture(str(img), PInches(1), PInches(1.8), width=PInches(2.5))
    slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = "Synthetic Metrics"
    data = ChartData(); data.categories = ["Q1", "Q2", "Q3", "Q4"]; data.add_series("Operations", (12, 19, 17, 24)); data.add_series("Quality", (7, 10, 11, 14))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, PInches(1), PInches(1.7), PInches(8), PInches(4.5), data)
    slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = "Risk Register"
    table = slide.shapes.add_table(4, 3, PInches(0.8), PInches(1.7), PInches(8.5), PInches(3.0)).table
    for c, text in enumerate(["Risk", "Owner", "Status"]): table.cell(0, c).text = text
    for r, row in enumerate([("Unexpected egress", "Security", "Blocked"), ("Data loss", "Operations", "Mitigated"), ("Model overreach", "Product", "AI OFF")], start=1):
        for c, text in enumerate(row): table.cell(r, c).text = text
    slide = prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text = "Decision"; slide.placeholders[1].text = "Preserve user data\nVerify before claiming success"
    for s in prs.slides:
        try:
            add_transition(s)
        except Exception:
            pass
        try:
            notes = s.notes_slide.notes_text_frame
            notes.text = "SYNTHETIC TEST DATA - speaker notes for compatibility testing."
        except Exception:
            pass
    prs.save(path)
    return path


def make_pdf() -> Path:
    path = OUT / "InstitutionalPacket.pdf"
    c = canvas.Canvas(str(path), pagesize=A4)
    for i in range(1, 5):
        c.setFont("Helvetica-Bold", 16); c.drawString(72, 780, "KHZ Institutional Packet")
        c.setFont("Helvetica", 11); c.drawString(72, 752, "SYNTHETIC TEST DATA")
        c.drawString(72, 724, f"Page {i} - local PDF viewing/search fixture")
        for line in range(12): c.drawString(72, 690 - line * 22, f"Synthetic record {i:02d}-{line+1:02d}: neutral test content for offline verification.")
        c.showPage()
    c.save(); return path


def main() -> int:
    img = OUT / "synthetic-image.png"; tiny_png(img)
    paths = [make_docx(img), make_xlsx(), make_pptx(img), make_pdf()]
    print("Generated:")
    for p in paths: print(p)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
